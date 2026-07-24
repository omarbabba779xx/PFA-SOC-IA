from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

SS = r"C:\Users\nextlevel\Desktop\PFA-SOC-IA\docs\evidence\final\PFA-FINAL-20260718-214637\presentation_finale\screenshots"

BG = RGBColor(0x0B, 0x0F, 0x14)
FG = RGBColor(0xF2, 0xF4, 0xF7)
ACCENT = RGBColor(0x3D, 0xB8, 0xFF)
MUTED = RGBColor(0x9A, 0xA5, 0xB1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = FG
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(15)
        p2.font.color.rgb = ACCENT


def add_image_slide(title, sub, img_file, caption=None):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, title, sub)
    img_path = f"{SS}\\{img_file}"
    left = Inches(0.6)
    top = Inches(1.25)
    max_w = Inches(12.13)
    max_h = Inches(5.75)
    from PIL import Image as PILImage
    im = PILImage.open(img_path)
    ratio = im.width / im.height
    w = max_w
    h = Emu(int(w / ratio))
    if h > max_h:
        h = max_h
        w = Emu(int(h * ratio))
    left = Emu(int((prs.slide_width - w) / 2))
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.4))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = MUTED
    return slide


def add_text_slide(title, sub, lines):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, title, sub)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = FG
        p.space_after = Pt(14)
    return slide


# 1. Title
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SOC Assisté par Intelligence Artificielle"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = FG
p2 = tf.add_paragraph()
p2.text = "Détection Wazuh -> Triage IA (Gemma2) -> TheHive -> Cortex -> MISP -> Notification"
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT
p2.space_before = Pt(16)
p3 = tf.add_paragraph()
p3.text = "Omar Babba - PFA EMSI Tanger - RUN_ID PFA-FINAL-20260718-214637"
p3.font.size = Pt(14)
p3.font.color.rgb = MUTED
p3.space_before = Pt(10)

# 2. Architecture / chaine logique (text only, minimal)
add_text_slide(
    "Chaine logique du pipeline",
    "De la detection a la notification - execution reelle, VM VirtualBox",
    [
        "ETAPE 1 - Wazuh (auditd) detecte une activite suspecte sur l'agent surveille",
        "ETAPE 2 - L'alerte reelle declenche le workflow Shuffle via webhook",
        "ETAPE 3 - Gemma2 9B (Ollama, local) effectue le triage : type, MITRE, criticite",
        "ETAPE 4 - TheHive cree automatiquement un cas avec le triage IA",
        "ETAPE 5 - Cortex enrichit l'IOC (ex: AbuseIPDB)",
        "ETAPE 6 - Selon la severite : evenement MISP (haute) ou tag simple (basse)",
        "ETAPE 7 - Notification finale envoyee",
        "Chaque etape critique est protegee par une garde de statut HTTP.",
    ],
)

# 3. ETAPE 1 - Wazuh overview
add_image_slide(
    "ETAPE 1 / 7 - Wazuh detecte",
    "Dashboard reel : agent actif, alertes classees par severite",
    "01_wazuh_overview.png",
)

# 4. ETAPE 1 - Wazuh alert (suite)
add_image_slide(
    "ETAPE 1 / 7 - Alerte reelle declenchee",
    "Regle 100103 - C2 beaconing - MITRE T1071 - agent soc-lab",
    "02_wazuh_alerte_c2_beaconing.png",
)

# 5. ETAPE 2-3 - Shuffle canvas (orchestration + triage)
add_image_slide(
    "ETAPE 2-3 / 7 - Shuffle recoit l'alerte et orchestre",
    "Webhook -> Gemma2 (triage) -> TheHive -> Cortex -> MISP/tag -> notification",
    "03_shuffle_canvas_complet.png",
)

# 6. ETAPE 3 - Gemma2 triage config
add_image_slide(
    "ETAPE 3 / 7 - Triage par IA (Gemma2 9B)",
    "Prompt reel envoye au modele local Ollama pour classification MITRE",
    "04_shuffle_config_gemma2.png",
)

# 7. ETAPE 4 - TheHive
add_image_slide(
    "ETAPE 4 / 7 - TheHive cree le cas",
    "Triage brut de Gemma2 visible dans la description du cas reel",
    "10_thehive_case22_gemma_triage.png",
)

# 7bis. Preuve de l'automatisation reelle (execution Shuffle correlee au cas TheHive)
add_image_slide(
    "Preuve : le cas TheHive vient bien de Shuffle",
    "Reponse HTTP reelle 201 de l'execution Shuffle - meme case id ~28720 / #22",
    "27_shuffle_execution_thehive_case_proof.png",
    caption="Correle avec la capture precedente : createdBy soc-pipeline52@thehive.local (compte de service), pas un humain.",
)

# 8. ETAPE 5 - Cortex jobs
add_image_slide(
    "ETAPE 5 / 7 - Cortex enrichit l'IOC",
    "Historique des analyses AbuseIPDB reelles, toutes en Success",
    "11_cortex_jobs_history.png",
)

# 9. ETAPE 5 - Cortex report (suite)
add_image_slide(
    "ETAPE 5 / 7 - Rapport d'enrichissement detaille",
    "Resultat reel de l'API AbuseIPDB sur l'IP suspecte",
    "12_cortex_job_report_abuseipdb.png",
)

# 10. ETAPE 6 - MISP event
add_image_slide(
    "ETAPE 6 / 7 - MISP partage la menace",
    "Evenement cree automatiquement par le pipeline (severite haute)",
    "13_misp_event11_header.png",
)

# 11. ETAPE 6 - MISP IOC (suite)
add_image_slide(
    "ETAPE 6 / 7 - IOC reel extrait et correle",
    "Attribut ip-dst correle a 4 autres evenements, IDS actif",
    "14_misp_event11_attribute_ioc.png",
)

# 12. ETAPE 7 - Notification
add_image_slide(
    "ETAPE 7 / 7 - Notification finale",
    "Dernier noeud du pipeline : alerte envoyee au canal de notification",
    "08_shuffle_config_notification.png",
)

# 11. Closing / results
add_text_slide(
    "Bilan",
    "Tests reels effectues le jour meme",
    [
        "3 executions completes reelles du pipeline aujourd'hui",
        "Run 100% vert : Gemma2 -> TheHive -> Cortex -> MISP -> notification",
        "Gardes d'echec verifiees en conditions reelles (TheHive 401, Cortex ES down)",
        "Aucune donnee simulee : alertes, cas, jobs et evenements sont reels",
    ],
)

out = r"C:\Users\nextlevel\Desktop\PFA-SOC-IA\docs\evidence\final\PFA-FINAL-20260718-214637\presentation_finale\PFA_SOC_IA_Presentation.pptx"
prs.save(out)
print("saved", out)
